import streamlit as st
import google.generativeai as genai
from docx import Document
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import time
import json
import os
import pandas as pd

# --- 1. DATA REPOSITORY ---
PROFICIENT_DEFS = {
    "Indicator 1": "Apply knowledge of content within and across curriculum teaching areas.",
    "Indicator 2": "Use a range of teaching strategies that enhance learner achievement in literacy and numeracy skills.",
    "Indicator 3": "Apply a range of teaching strategies to develop critical and creative thinking, as well as other higher-order thinking skills.",
    "Indicator 4": "Manage classroom structure to engage learners in meaningful exploration and discovery.",
    "Indicator 5": "Manage learner behavior constructively by applying positive and non-violent discipline.",
    "Indicator 6": "Use differentiated, developmentally appropriate learning experiences to address learner needs.",
    "Indicator 7": "Plan, manage and implement developmentally sequenced teaching and learning processes.",
    "Indicator 8": "Select, develop, organize and use appropriate teaching and learning resources, including ICT.",
    "Indicator 9": "Design, select, organize and use diagnostic, formative and summative assessment strategies."
}

HIGHLY_PROFICIENT_DEFS = {
    "Indicator 1": "Model effective applications of content knowledge within and across curriculum areas.",
    "Indicator 2": "Collaborate with colleagues to enhance learner achievement in literacy and numeracy.",
    "Indicator 3": "Develop and apply effective communication strategies to support learner participation.",
    "Indicator 4": "Work with colleagues to provide safe and secure learning environments.",
    "Indicator 5": "Exhibit effective practices in the management of learner behavior.",
    "Indicator 6": "Work with colleagues to share differentiated, developmentally appropriate learning experiences.",
    "Indicator 7": "Develop and apply effective strategies in the planning and management of sequences of teaching and learning.",
    "Indicator 8": "Advise and guide colleagues in the selection and use of appropriate resources, including ICT.",
    "Indicator 9": "Work with colleagues to design, select, organize and use diagnostic and summative assessments."
}

# --- 2. THEMED EXPORT ENGINES ---
def convert_to_pptx(text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Instructional Presentation"
    slide.placeholders[1].text = f"Context: {st.session_state.get('context', 'DepEd')}\nDeveloped by Eleazer A. Meriño"
    sections = text.split("###")
    for section in sections:
        if section.strip():
            lines = section.strip().split("\n")
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = lines[0].strip()
            slide.placeholders[1].text = "\n".join(lines[1:])[:1000]
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

def convert_to_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=10)
    clean_text = text.encode('latin-1', 'ignore').decode('latin-1')
    pdf.multi_cell(0, 7, txt=clean_text)
    return bytes(pdf.output())

def convert_to_docx(text):
    doc = Document()
    doc.add_heading('Jargon Buster Pro - Instructional Suite', 0)
    clean_text = "".join(c for c in text if ord(c) > 31 or c in '\n\r\t')
    doc.add_paragraph(clean_text)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# --- 3. PREMIUM UI CONFIGURATION ---
st.set_page_config(page_title="Jargon Buster Pro | Commercial", page_icon="🚀", layout="wide")

st.markdown("""
    <style>
    .stApp { background-color: #0A0C10; color: #E6EDF3; }
    [data-testid="stSidebar"] { background-color: #111418 !important; border-right: 1px solid #30363D; }
    .stButton>button { 
        border-radius: 8px; 
        background: linear-gradient(90deg, #1E40AF 0%, #3B82F6 100%); 
        color: white; border: none; font-weight: 700; height: 3em; width: 100%;
    }
    .stButton>button:hover { background: linear-gradient(90deg, #3B82F6 0%, #60A5FA 100%); transform: translateY(-1px); }
    div[data-testid="stExpander"] { background-color: #161B22; border: 1px solid #30363D; border-radius: 10px; }
    .license-badge { padding: 5px 10px; border-radius: 20px; font-size: 12px; font-weight: bold; text-align: center; margin-bottom: 10px; }
    </style>
    """, unsafe_allow_html=True)

if "current_content" not in st.session_state: 
    st.session_state.current_content = ""

# --- 4. SIDEBAR SETUP & GUIDE ---
with st.sidebar:
    st.markdown("<h2 style='color: #3B82F6;'>🚀 JB PRO</h2>", unsafe_allow_html=True)
    st.caption("Commercial Suite v9.0 | The Ultimate Edition")
    
    api_key = st.text_input("🔑 System Activation Key (Gemini API)", type="password")
    if api_key:
        st.markdown('<div class="license-badge" style="background-color: #238636; color: white;">LICENSE ACTIVE</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="license-badge" style="background-color: #DA3633; color: white;">KEY REQUIRED</div>', unsafe_allow_html=True)
        with st.expander("📖 Setup Wizard: How to start"):
            st.markdown("""
            **Step 1:** Visit [Google AI Studio](https://aistudio.google.com/app/apikey).
            **Step 2:** Create an API key.
            **Step 3:** Paste it above to activate the AI engine!
            """)
    
    st.divider()
    nav_mode = st.radio("🏠 Hub", ["Teacher Workspace", "Principal Dashboard"])
    rank = st.selectbox("Your Rank", ["Proficient (T1-T3)", "Highly Proficient (MT1-MT4)"])
    
    st.divider()
    st.markdown(f"<div style='text-align: center; color: #8B949E; font-size: 11px;'><strong>© 2026 Eleazer A. Meriño</strong><br>Pangasinan Division</div>", unsafe_allow_html=True)

# --- 5. TEACHER WORKSPACE ---
if nav_mode == "Teacher Workspace":
    st.title("Instructional Strategy Room")
    
    if not api_key:
        st.info("👋 **Welcome to Jargon Buster Pro.** Please activate your system using an API key in the sidebar.")
    else:
        is_hp = "Highly Proficient" in rank
        curr_defs = HIGHLY_PROFICIENT_DEFS if is_hp else PROFICIENT_DEFS
        
        c_in, c_out = st.columns([1, 1.4], gap="large")

        with c_in:
            with st.container(border=True):
                st.markdown("### ⚙️ Core Settings")
                
                # NEW: Feature 4 - Weekly Matrix Mode
                lesson_format = st.radio("📅 Lesson Format", ["Single Lesson (5E Plan)", "Weekly Matrix (DepEd DLL)"], horizontal=True)
                
                lang = st.selectbox("Output Language", ["English", "Tagalog", "Ilokano", "Pangasinan"])
                context = st.text_input("Local Context", value="Binmaley / Binalonan")
                st.session_state['context'] = context
                
                selected = st.multiselect("🎯 Target Indicators (RPMS-PPST):", list(curr_defs.keys()))
                
                st.divider()
                # NEW: Feature 3 - Assessment Generator
                st.markdown("### 📝 Extra Resources")
                include_quiz = st.checkbox("➕ Generate 10-Item MCQ Quiz + Answer Key")
                include_las = st.checkbox("➕ Generate Learning Activity Sheet (LAS)", value=True)
                include_rubric = st.checkbox("➕ Generate Analytical Rubric", value=True)
                
                st.divider()
                melc = st.text_area("Paste MELC / Topic:", height=120, placeholder="e.g., Grade 8 Science: Reflections of Light and Laws of Motion")

            if st.button("Generate Master Suite ✨", use_container_width=True):
                if not melc: 
                    st.error("Please input a lesson topic.")
                else:
                    # NEW: Feature 1 - Real-Time Streaming Setup
                    try:
                        genai.configure(api_key=api_key)
                        model = genai.GenerativeModel('gemini-2.5-flash')
                        ind_txt = "\n".join([f"{k}: {curr_defs[k]}" for k in selected])
                        
                        # Dynamic Prompt Construction based on new features
                        format_instruction = "Generate a standard 5E Lesson Plan." if lesson_format == "Single Lesson (5E Plan)" else "Generate a 5-day Daily Lesson Log (DLL) matrix covering Day 1 to Day 5."
                        extras = []
                        if include_las: extras.append("a student Learning Activity Sheet (LAS)")
                        if include_rubric: extras.append("an analytical rubric")
                        if include_quiz: extras.append("a 10-item multiple-choice summative assessment with a separate Answer Key")
                        
                        extras_instruction = f"Also generate: {', '.join(extras)}." if extras else ""
                        
                        prompt = f"""
                        Role: DepEd {rank}. Lang: {lang}. Context: {context}. 
                        Indicators to hit: {ind_txt}. 
                        Topic: {melc}. 
                        
                        Instructions:
                        {format_instruction}
                        {extras_instruction}
                        
                        Mark specific moves with [HIT: INDICATOR X]. 
                        Use '###' for major headers to support the PowerPoint parser.
                        """
                        
                        with c_out:
                            st.markdown("### 🧠 AI Generation in Progress...")
                            output_placeholder = st.empty()
                            full_response = ""
                            
                            response = model.generate_content(prompt, stream=True)
                            
                            for chunk in response:
                                full_response += chunk.text
                                # Typing cursor effect
                                output_placeholder.markdown(full_response + "▌")
                            
                            output_placeholder.empty() # Clear placeholder once done
                            st.session_state.current_content = full_response
                            st.rerun() # Refresh to show the editor

                    except Exception as e: 
                        st.error(f"API Error: {e}")

        with c_out:
            if st.session_state.current_content:
                st.markdown("### 📝 Human-in-the-Loop Editor")
                st.caption("AI gets you 90% there. Use this editor to tweak the final 10% before exporting.")
                
                # NEW: Feature 2 - The In-App Editor
                edited_text = st.text_area(
                    "Fine-tune your document:", 
                    value=st.session_state.current_content, 
                    height=500,
                    label_visibility="collapsed"
                )
                
                st.divider()
                st.markdown("### 📥 Commercial Export")
                # The export buttons now pull directly from the `edited_text` variable!
                ea, eb, ec = st.columns(3)
                ea.download_button("📂 PDF Suite", convert_to_pdf(edited_text), "Instructional_Suite.pdf", use_container_width=True)
                eb.download_button("📂 Word Doc", convert_to_docx(edited_text), "DepEd_Plan.docx", use_container_width=True)
                ec.download_button("📂 PPT Slides", convert_to_pptx(edited_text), "Presentation.pptx", use_container_width=True)

# --- 6. PRINCIPAL DASHBOARD ---
else:
    st.title("📊 Principal Analytics")
    st.info("System tracking for localized usage is active.")